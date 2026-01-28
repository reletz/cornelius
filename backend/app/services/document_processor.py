"""
Document processing service - OCR and text extraction.
"""
import asyncio
from pathlib import Path
from typing import Optional
import logging

from app.core.config import settings

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """
    Handles document text extraction using tiered strategy:
    - Tier 1: Direct extraction (python-pptx, pypdfium2, python-docx)
    - Tier 2: OCR fallback (Tesseract)
    """
    
    SUPPORTED_EXTENSIONS = {
        '.pptx': 'pptx',
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.png': 'image',
        '.jpg': 'image',
        '.jpeg': 'image',
    }
    
    async def extract_text(self, file_path: str) -> str:
        """
        Extract text from document using appropriate method.
        """
        path = Path(file_path)
        extension = path.suffix.lower()
        
        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {extension}")
        
        doc_type = self.SUPPORTED_EXTENSIONS[extension]
        
        # Run extraction in thread pool to avoid blocking
        loop = asyncio.get_event_loop()
        
        try:
            if doc_type == 'pptx':
                text = await loop.run_in_executor(None, self._extract_pptx, file_path)
            elif doc_type == 'pdf':
                text = await loop.run_in_executor(None, self._extract_pdf, file_path)
            elif doc_type == 'docx':
                text = await loop.run_in_executor(None, self._extract_docx, file_path)
            elif doc_type == 'image':
                text = await loop.run_in_executor(None, self._extract_image, file_path)
            else:
                raise ValueError(f"Unknown document type: {doc_type}")
            
            # Check if we got meaningful text
            if not text or len(text.strip()) < 50:
                logger.warning(f"Insufficient text from direct extraction, trying OCR: {file_path}")
                if doc_type == 'pdf':
                    text = await loop.run_in_executor(None, self._ocr_pdf, file_path)
            
            return text.strip()
            
        except asyncio.TimeoutError:
            raise TimeoutError(f"Processing timed out for: {file_path}")
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            raise
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint file."""
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"\n--- Slide {slide_num} ---\n"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                    
                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text for cell in row.cells if cell.text]
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            text_parts.extend(slide_text)
        
        return "\n".join(text_parts)
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF using pypdfium2."""
        import pypdfium2 as pdfium
        
        pdf = pdfium.PdfDocument(file_path)
        text_parts = []
        
        for page_num, page in enumerate(pdf, 1):
            text_parts.append(f"\n--- Page {page_num} ---\n")
            textpage = page.get_textpage()
            text_parts.append(textpage.get_text_range())
        
        return "\n".join(text_parts)
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from Word document."""
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Handle tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells if cell.text]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n".join(text_parts)
    
    def _extract_image(self, file_path: str) -> str:
        """Extract text from image using Tesseract OCR."""
        import pytesseract
        from PIL import Image
        
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang='eng+ind')
        return text
    
    def _ocr_pdf(self, file_path: str) -> str:
        """
        OCR fallback for scanned PDFs.
        Converts PDF pages to images and runs Tesseract.
        """
        import pypdfium2 as pdfium
        import pytesseract
        from PIL import Image
        import io
        
        pdf = pdfium.PdfDocument(file_path)
        text_parts = []
        
        for page_num, page in enumerate(pdf, 1):
            text_parts.append(f"\n--- Page {page_num} (OCR) ---\n")
            
            # Render page to image
            bitmap = page.render(scale=2.0)  # Higher scale for better OCR
            pil_image = bitmap.to_pil()
            
            # Run OCR
            text = pytesseract.image_to_string(pil_image, lang='eng+ind')
            text_parts.append(text)
        
        return "\n".join(text_parts)


# Singleton instance
document_processor = DocumentProcessor()
